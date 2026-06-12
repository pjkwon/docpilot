from __future__ import annotations

import io
import os
import zipfile
from pathlib import Path

import pytest

# Minimal HWPX content.hml XML with one heading and one placeholder
_HWPX_CONTENT = """<?xml version="1.0" encoding="UTF-8"?>
<hml xmlns:hp="http://www.hancom.co.kr/hwpml/2012/paragraph"
     xmlns:hc="http://www.hancom.co.kr/hwpml/2012/core">
  <hp:body>
    <hp:sec>
      <hp:p>
        <hp:pPr>
          <hp:rPr>
            <hp:sz hp:val="2000"/>
            <hp:b hp:val="true"/>
          </hp:rPr>
        </hp:pPr>
        <hp:t>서론</hp:t>
      </hp:p>
      <hp:p>
        <hp:t>{{서론}}</hp:t>
      </hp:p>
      <hp:p>
        <hp:pPr>
          <hp:rPr>
            <hp:sz hp:val="1800"/>
            <hp:b hp:val="true"/>
          </hp:rPr>
        </hp:pPr>
        <hp:t>결론</hp:t>
      </hp:p>
      <hp:p>
        <hp:t>{{결론}}</hp:t>
      </hp:p>
    </hp:sec>
  </hp:body>
</hml>""".encode("utf-8")


def _make_hwpx(path: Path, content_xml: bytes = _HWPX_CONTENT) -> Path:
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("mimetype", "application/hwp+zip", compress_type=zipfile.ZIP_STORED)
        zf.writestr("Contents/content.hml", content_xml)
        zf.writestr("META-INF/container.xml", b"<container/>")
    return path


@pytest.fixture()
def sample_hwpx(tmp_path: Path) -> Path:
    return _make_hwpx(tmp_path / "sample.hwpx")


@pytest.fixture()
def sample_txt(tmp_path: Path) -> Path:
    p = tmp_path / "sample.txt"
    p.write_text("2025년 사업 계획서\n\n핵심 목표: 매출 30% 성장\n세부 계획: ...", encoding="utf-8")
    return p


@pytest.fixture()
def sample_pptx(tmp_path: Path) -> Path:
    pptx = pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "2025 사업 계획"
    slide.placeholders[1].text = "핵심 내용"

    path = tmp_path / "sample.pptx"
    prs.save(str(path))
    return path


@pytest.fixture()
def hwpx_template(tmp_path: Path) -> Path:
    return _make_hwpx(tmp_path / "template.hwpx")


@pytest.fixture()
def make_hwpx(tmp_path: Path):
    """Factory fixture: make_hwpx(name, xml=...) -> Path"""
    def _factory(name: str = "doc.hwpx", content_xml: bytes = _HWPX_CONTENT) -> Path:
        return _make_hwpx(tmp_path / name, content_xml)
    return _factory


# ---------------------------------------------------------------------------
# DocPilot instance fixtures
# ---------------------------------------------------------------------------

# LLM이 필요한 테스트는 실제 키가 없으면 자동 skip
requires_llm = pytest.mark.skipif(
    not os.environ.get("ANTHROPIC_API_KEY"),
    reason="ANTHROPIC_API_KEY not set",
)


@pytest.fixture(scope="session")
def embed_fn():
    """BGE-M3 임베딩 함수 (1024차원, EMBEDDING_DIM과 일치). 세션 전체에서 모델을 한 번만 로드."""
    from docpilot.search.embedding import bge_embed_fn
    return bge_embed_fn()


@pytest.fixture()
def pilot(tmp_path: Path):
    """
    LLM 호출 없이 인덱싱·검색만 테스트할 때 사용.
    테스트마다 격리된 DB를 생성하므로 side effect 없음.
    """
    from docpilot import DocPilot
    return DocPilot(
        api_key="sk-test",  # init 검증 통과용 더미 키 — map() 호출 시엔 실제 키 필요
        database_url=f"sqlite:///{tmp_path / 'test.db'}",
    )


@pytest.fixture()
def pilot_with_embed(tmp_path: Path, embed_fn):
    """
    임베딩 포함 인덱싱·벡터 검색 테스트용.
    embed_fn은 session-scoped라 모델 로딩 비용은 1회만 발생.
    """
    from docpilot import DocPilot
    return DocPilot(
        api_key="sk-test",
        database_url=f"sqlite:///{tmp_path / 'test.db'}",
        embed_fn=embed_fn,
    )


@pytest.fixture()
def pilot_llm(tmp_path: Path):
    """
    실제 LLM 호출이 필요한 통합 테스트용.
    ANTHROPIC_API_KEY 없으면 테스트가 자동 skip되지 않으므로,
    테스트 함수에 @requires_llm 마커를 함께 붙여서 사용.
    """
    from docpilot import DocPilot
    return DocPilot(
        database_url=f"sqlite:///{tmp_path / 'test.db'}",
    )
